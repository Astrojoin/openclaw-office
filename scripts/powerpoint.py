#!/usr/bin/env python3
"""CRUD operations for .pptx files — purely offline, no network calls.

Works with bytes in memory (primary) or local file paths (optional).
Uses python-pptx for all presentation manipulation.

Usage:
  from powerpoint import PptxDoc

  # From bytes (downloaded from OneDrive, etc.)
  pptx = PptxDoc(data=bytes_data)
  slides = pptx.read()
  pptx.update({"operations": [{"method": "add_slide", "args": ["Title Slide"]}]})
  new_bytes = pptx.to_bytes()

  # From local file
  pptx = PptxDoc(path="/tmp/deck.pptx")

  # Create blank presentation
  pptx = PptxDoc.create()
  pptx.update({"operations": [{"method": "add_slide", "kwargs": {"layout": "Title Slide", "title": "My Deck"}}])
"""

from __future__ import annotations

import io
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Layout name → index mapping helper
_LAYOUT_MAP = {
    "title slide": 0,
    "title and content": 1,
    "section header": 2,
    "two content": 3,
    "comparison": 4,
    "title only": 5,
    "blank": 6,
}


class PptxDoc:
    """In-memory .pptx editor using python-pptx."""

    def __init__(self, data: bytes | None = None, path: str | None = None):
        self.path = Path(path) if path else None
        self._prs: Presentation | None = None

        if data is not None:
            self._prs = Presentation(io.BytesIO(data))
        elif path is not None:
            self._prs = Presentation(path)
        else:
            raise ValueError("Provide either 'data' (bytes) or 'path' (str), or use PptxDoc.create()")

    @classmethod
    def create(cls) -> "PptxDoc":
        """Create a blank presentation."""
        instance = cls.__new__(cls)
        instance.path = None
        instance._prs = Presentation()
        return instance

    @property
    def prs(self) -> Presentation:
        return self._prs

    def _resolve_layout(self, layout_name: str):
        """Resolve a layout by name or index. Returns slide_layout object."""
        name_lower = layout_name.lower()
        for sl in self._prs.slide_layouts:
            if sl.name.lower() == name_lower:
                return sl
        # Fallback: try numeric index
        idx = _LAYOUT_MAP.get(name_lower, 0)
        if idx < len(self._prs.slide_layouts):
            return self._prs.slide_layouts[idx]
        return self._prs.slide_layouts[0]

    # ── READ ────────────────────────────────────────────────────────────

    def read(self, mode: str = "plain") -> str:
        """Read presentation content.

        Args:
            mode: "plain" for readable text, "structured" for JSON
        """
        if mode == "structured":
            return self._read_structured()
        return self._read_plain()

    def _read_plain(self) -> str:
        """Extract all slide text as readable string."""
        lines = []
        for i, slide in enumerate(self._prs.slides, 1):
            lines.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(f"  {text}")
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text for cell in row.cells]
                        lines.append("  | " + " | ".join(cells) + " |")
        return "\n".join(lines)

    def _read_structured(self) -> str:
        """Extract presentation as structured JSON."""
        slides = []
        for i, slide in enumerate(self._prs.slides, 1):
            shapes = []
            for shape in slide.shapes:
                s = {"name": shape.name, "shape_type": str(shape.shape_type)}
                if shape.has_text_frame:
                    paragraphs = []
                    for para in shape.text_frame.paragraphs:
                        paragraphs.append({
                            "text": para.text,
                            "level": para.level,
                        })
                    s["text_frame"] = paragraphs
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text for cell in row.cells])
                    s["table"] = rows
                shapes.append(s)
            slides.append({"slide_number": i, "shapes": shapes})
        return json.dumps({"slides": slides}, indent=2, ensure_ascii=False)

    # ── UPDATE ──────────────────────────────────────────────────────────

    def update(self, spec: dict) -> str:
        """Apply operations to the presentation.

        Args:
            spec: {"operations": [{"method": "...", "args": [...], "kwargs": {...}}]}
        """
        operations = spec.get("operations", [])
        results = []
        for op in operations:
            method = op.get("method", "")
            args = op.get("args", [])
            kwargs = op.get("kwargs", {})
            result = self._apply_operation(method, args, kwargs)
            results.append(result)
        return f"Applied {len(results)} operations: " + "; ".join(results)

    def _apply_operation(self, method: str, args: list, kwargs: dict) -> str:
        prs = self._prs

        # ── Add slide ───────────────────────────────────────────────
        if method == "add_slide":
            layout_name = args[0] if args else kwargs.get("layout", "Title and Content")
            layout = self._resolve_layout(layout_name)
            slide = prs.slides.add_slide(layout)

            # Set title if provided
            title = kwargs.get("title", None)
            if title and slide.shapes.title:
                slide.shapes.title.text = title

            # Set subtitle/body if provided
            body_text = kwargs.get("body", kwargs.get("subtitle", None))
            if body_text:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # Subtitle/body
                        shape.text = body_text
                        break
            return f"add_slide(layout='{layout_name}')"

        # ── Remove slide ────────────────────────────────────────────
        if method == "remove_slide":
            idx = args[0] if args else kwargs.get("index", 0)
            if 0 <= idx < len(prs.slides):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]
                return f"remove_slide({idx})"
            return f"ERROR: slide index {idx} out of range"

        # ── Add text to slide ───────────────────────────────────────
        if method == "add_textbox":
            slide_idx = kwargs.pop("slide_index", len(prs.slides) - 1)
            slide = prs.slides[slide_idx]
            left = kwargs.get("left", Inches(1))
            top = kwargs.get("top", Inches(2))
            width = kwargs.get("width", Inches(8))
            height = kwargs.get("height", Inches(1))
            text = args[0] if args else kwargs.get("text", "")

            # Convert numeric values to EMU if they look like inches
            if isinstance(left, (int, float)) and left < 20:
                left = Inches(left)
            if isinstance(top, (int, float)) and top < 20:
                top = Inches(top)
            if isinstance(width, (int, float)) and width < 20:
                width = Inches(width)
            if isinstance(height, (int, float)) and height < 20:
                height = Inches(height)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = text

            if "font_size" in kwargs:
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(kwargs["font_size"])
            if kwargs.get("bold"):
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.bold = True
            return f"add_textbox(slide={slide_idx}, text='{text[:40]}')"

        # ── Add table to slide ──────────────────────────────────────
        if method == "add_table":
            slide_idx = kwargs.pop("slide_index", len(prs.slides) - 1)
            slide = prs.slides[slide_idx]
            rows = args[0] if len(args) > 0 else kwargs.get("rows", 2)
            cols = args[1] if len(args) > 1 else kwargs.get("cols", 2)
            left = kwargs.get("left", Inches(1))
            top = kwargs.get("top", Inches(2))
            width = kwargs.get("width", Inches(8))
            height = kwargs.get("height", Inches(3))

            if isinstance(left, (int, float)) and left < 20:
                left = Inches(left)
            if isinstance(top, (int, float)) and top < 20:
                top = Inches(top)

            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            if "data" in kwargs:
                for i, row_data in enumerate(kwargs["data"]):
                    if i < len(table.rows):
                        for j, val in enumerate(row_data):
                            if j < len(table.columns):
                                table.cell(i, j).text = str(val)
            return f"add_table(slide={slide_idx}, {rows}x{cols})"

        # ── Set cell text ───────────────────────────────────────────
        if method == "set_table_cell":
            slide_idx = kwargs.pop("slide_index", len(prs.slides) - 1)
            slide = prs.slides[slide_idx]
            table_idx = kwargs.get("table_index", 0)
            row_idx = args[0] if len(args) > 0 else kwargs.get("row", 0)
            col_idx = args[1] if len(args) > 1 else kwargs.get("col", 0)
            text = args[2] if len(args) > 2 else kwargs.get("text", "")

            shape = None
            tcount = 0
            for s in slide.shapes:
                if s.has_table:
                    if tcount == table_idx:
                        shape = s
                        break
                    tcount += 1
            if shape:
                shape.table.cell(row_idx, col_idx).text = text
                return f"set_table_cell(slide={slide_idx},table={table_idx},[{row_idx},{col_idx}])"
            return f"ERROR: no table {table_idx} on slide {slide_idx}"

        # ── Add image to slide ──────────────────────────────────────
        if method == "add_picture":
            slide_idx = kwargs.pop("slide_index", len(prs.slides) - 1)
            slide = prs.slides[slide_idx]
            img_path = args[0] if args else kwargs.get("path", "")
            left = kwargs.get("left", Inches(1))
            top = kwargs.get("top", Inches(2))
            width = kwargs.get("width", None)
            height = kwargs.get("height", None)

            if isinstance(left, (int, float)) and left < 20:
                left = Inches(left)
            if isinstance(top, (int, float)) and top < 20:
                top = Inches(top)

            w = Inches(width) if width and width < 20 else width
            h = Inches(height) if height and height < 20 else height

            if w and h:
                slide.shapes.add_picture(img_path, left, top, w, h)
            elif w:
                slide.shapes.add_picture(img_path, left, top, width=w)
            else:
                slide.shapes.add_picture(img_path, left, top)
            return f"add_picture(slide={slide_idx}, '{img_path}')"

        # ── Replace text in slide ───────────────────────────────────
        if method == "replace_text":
            old = args[0] if len(args) > 0 else kwargs.get("old", "")
            new = args[1] if len(args) > 1 else kwargs.get("new", "")
            count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if old in run.text:
                                    run.text = run.text.replace(old, new)
                                    count += 1
            return f"replace_text: {count} replacements"

        # ── Add speaker notes ───────────────────────────────────────
        if method == "add_notes":
            slide_idx = args[0] if args else kwargs.get("slide_index", len(prs.slides) - 1)
            text = args[1] if len(args) > 1 else kwargs.get("text", "")
            slide = prs.slides[slide_idx]
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = text
            return f"add_notes(slide={slide_idx})"

        # Fallback
        return f"ERROR: unknown method '{method}'"

    # ── DELETE (local only) ─────────────────────────────────────────────

    def delete(self) -> str:
        if self.path and self.path.exists():
            self.path.unlink()
            return f"Deleted: {self.path}"
        return "No local file to delete (in-memory document). Use onedrive.py for cloud deletion."

    # ── OUTPUT ──────────────────────────────────────────────────────────

    def to_bytes(self) -> bytes:
        buf = io.BytesIO()
        self._prs.save(buf)
        buf.seek(0)
        return buf.read()

    def save(self, path: str | None = None) -> str:
        target = Path(path) if path else self.path
        if target is None:
            raise ValueError("No path specified. Use to_bytes() for in-memory output.")
        target.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(target))
        return f"Saved: {target}"


# ── CLI ─────────────────────────────────────────────────────────────────────

def main():
    pptx = PptxDoc.create()
    pptx.update({"operations": [
        {"method": "add_slide", "kwargs": {"layout": "Title Slide", "title": "OneDrive Skill Test"}},
        {"method": "add_slide", "kwargs": {"layout": "Title and Content", "title": "Slide 2", "body": "Content here"}},
    ]})
    print("=== Presentation content ===")
    print(pptx.read())
    print(f"\n=== Bytes size: {len(pptx.to_bytes())} bytes ===")


if __name__ == "__main__":
    main()
